import os, tempfile, shutil, re, logging, asyncio
from typing import List, Dict, Any
from fastapi import FastAPI, Request, UploadFile, File, Form, HTTPException
from fastapi.responses import HTMLResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
import httpx
import google.generativeai as genai

# --- Logging ---
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# --- App + UI ---
app = FastAPI(title="Markdown → PowerPoint")
templates = Jinja2Templates(directory="templates")

if os.path.isdir("static"):
    app.mount("/static", StaticFiles(directory="static"), name="static")

BUILTIN_TEMPLATES = {
    "minimal": os.path.join("templates_pack", "minimal.pptx"),
    "corporate": os.path.join("templates_pack", "corporate.pptx"),
}

# --- LLM Providers ---
LLM_PROVIDERS = {
    "openai": {
        "url": "https://api.openai.com/v1/chat/completions",
        "model": "gpt-4o-mini",
    },
    "anthropic": {
        "url": "https://api.anthropic.com/v1/messages",
        "model": "claude-3-sonnet-20240229",
        "version": "2023-06-01",
    },
    "gemini": {
        "model": "gemini-1.5-flash",
    },
    "aipipe": {
        "url": "https://aipipe.org/openrouter/v1/chat/completions",
        "model": "GPT-4o",
    },
}

# --- User-friendly error messages ---
USER_FRIENDLY_ERRORS = {
    "invalid_api_key": "Invalid API key. Please check your API key and try again.",
    "authentication_failed": "Authentication failed. Please verify your API credentials.",
    "insufficient_credits": "Your account has insufficient credits. Please add credits to your API account.",
    "credit_balance": "Your API account balance is too low. Please add credits to continue.",
    "billing": "Billing issue detected. Please check your account billing settings.",
    "payment_required": "Payment required. Please update your billing information.",
    "rate_limit": "Too many requests. Please wait a moment and try again.",
    "quota_exceeded": "API quota exceeded. Please try again later or upgrade your plan.",
    "timeout": "Request timed out. Please check your internet connection and try again.",
    "network_error": "Network connection error. Please check your internet connection.",
    "server_error": "Service temporarily unavailable. Please try again in a few moments.",
    "service_unavailable": "The API service is currently unavailable. Please try again later.",
    "invalid_request": "Invalid request. Please check your input and try again.",
    "permission_denied": "Permission denied. Please check your account permissions.",
    "default": "An unexpected error occurred. Please try again or contact support if the problem persists."
}

def get_user_friendly_error(error_text: str) -> str:
    error_text_lower = error_text.lower()
    if any(k in error_text_lower for k in ["invalid api key", "authentication", "unauthorized", "401", "403"]):
        return USER_FRIENDLY_ERRORS["invalid_api_key"]
    elif any(k in error_text_lower for k in ["credit", "balance", "billing", "payment", "insufficient"]):
        return USER_FRIENDLY_ERRORS["insufficient_credits"]
    elif any(k in error_text_lower for k in ["rate limit", "too many requests", "quota", "429"]):
        return USER_FRIENDLY_ERRORS["rate_limit"]
    elif any(k in error_text_lower for k in ["timeout", "timed out"]):
        return USER_FRIENDLY_ERRORS["timeout"]
    elif any(k in error_text_lower for k in ["network", "connection"]):
        return USER_FRIENDLY_ERRORS["network_error"]
    elif any(k in error_text_lower for k in ["server error", "service unavailable", "503", "502", "500"]):
        return USER_FRIENDLY_ERRORS["server_error"]
    elif any(k in error_text_lower for k in ["invalid request", "bad request", "400"]):
        return USER_FRIENDLY_ERRORS["invalid_request"]
    elif any(k in error_text_lower for k in ["permission", "forbidden"]):
        return USER_FRIENDLY_ERRORS["permission_denied"]
    return USER_FRIENDLY_ERRORS["default"]

@app.get("/", response_class=HTMLResponse)
def index(request: Request):
    return templates.TemplateResponse("index.html", {
        "request": request,
        "builtin_templates": BUILTIN_TEMPLATES.keys(),
        "llm_providers": list(LLM_PROVIDERS.keys()),
    })

# ---------------------------
# Cleaning helpers (safe)
# ---------------------------
MD_BOLD = re.compile(r"\*\*")  # remove ** used for markdown bold

def strip_leading_number(s: str) -> str:
    return re.sub(r"^\s*\d+[.)-]?\s*", "", s.strip())

def strip_slide_prefix(s: str) -> str:
    return re.sub(r"^\s*slide\s*\d*\s*:\s*", "", s, flags=re.IGNORECASE).strip()

def strip_bullet_marker(s: str) -> str:
    return re.sub(r"^\s*[-*•]\s+", "", s).strip()

def clean_inline(s: str) -> str:
    s = MD_BOLD.sub("", s)
    return s.strip()

# ---------------------------
# LLM Integration
# ---------------------------
async def call_llm(provider: str, api_key: str, prompt: str, guidance: str = "") -> str:
    if provider not in LLM_PROVIDERS:
        raise HTTPException(status_code=400, detail="Unsupported LLM provider")

    config = LLM_PROVIDERS[provider]

    system_prompt = """You are a presentation writing assistant.
Turn the input into a polished, ready-to-use PowerPoint deck.

Rules:
- Slide 1: Big centered title + optional subtitle.
- Slide 2 onwards: Title at top, with bullets.
- Each bullet = short headline + 1–2 sentence explanation.
- Headline should be bold, 5–10 words max.
- Explanation should be plain text under the headline.
- Use '•' for bullets (or ).
- No meta-text, no placeholders, no [square brackets].
- Do not output markdown headings like ##; just plain lines.

Format exactly like this:

Slide 1: Title
Subtitle: Optional

Slide 2: Slide Title
- Headline 1: Explanation sentence.
- Headline 2: Explanation sentence.
- Headline 3: Explanation sentence.
"""

    if guidance:
        system_prompt += f"\nAdditional guidance:\n{guidance}\n"

    try:
        if provider == "openai":
            headers = {"Content-Type": "application/json", "Authorization": f"Bearer {api_key}"}
            messages = [{"role": "system", "content": system_prompt},
                        {"role": "user", "content": prompt}]
            data = {"model": config["model"], "messages": messages, "temperature": 0.4}
            async with httpx.AsyncClient(timeout=60.0) as client:
                r = await client.post(config["url"], json=data, headers=headers)
                r.raise_for_status()
                return r.json()["choices"][0]["message"]["content"]

        elif provider == "anthropic":
            if not api_key.startswith("sk-ant-"):
                raise HTTPException(status_code=400, detail="Invalid Anthropic API key format. Should start with 'sk-ant-'")
            headers = {
                "Content-Type": "application/json",
                "x-api-key": api_key,
                "anthropic-version": config.get("version", "2023-06-01")
            }
            data = {
                "model": config["model"],
                "max_tokens": 4000,
                "messages": [{"role": "user", "content": f"{system_prompt}\n\n{prompt}"}],
                "temperature": 0.4
            }
            async with httpx.AsyncClient(timeout=60.0) as client:
                r = await client.post(config["url"], json=data, headers=headers)
                r.raise_for_status()
                res = r.json()
                return res["content"][0]["text"]

        elif provider == "gemini":
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel(
                model_name=config["model"],
                generation_config={"temperature": 0.3, "max_output_tokens": 4000},
                safety_settings=[{"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"}],
            )
            resp = model.generate_content(f"{system_prompt}\n\n{prompt}", request_options={"timeout": 60})
            return getattr(resp, "text", "").strip() or "Slide 1: Untitled\nSubtitle:"

        elif provider == "aipipe":
            headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
            data = {
                "model": config["model"],
                "messages": [{"role": "system", "content": system_prompt},
                             {"role": "user", "content": prompt}],
            }
            async with httpx.AsyncClient(timeout=60.0) as client:
                r = await client.post(config["url"], json=data, headers=headers)
                r.raise_for_status()
                result = r.json()
                return result.get("choices", [{}])[0].get("message", {}).get("content", "Slide 1: Untitled")

    except httpx.HTTPStatusError as e:
        error_detail = e.response.text if hasattr(e.response, 'text') else str(e)
        logger.error(f"{provider.upper()} API technical error: {error_detail}")
        user_message = get_user_friendly_error(error_detail)
        raise HTTPException(status_code=400, detail=user_message)

    except Exception as e:
        logger.error(f"{provider.upper()} request failed: {str(e)}")
        user_message = get_user_friendly_error(str(e))
        raise HTTPException(status_code=500, detail=user_message)

# ---------------------------
# Parsing LLM output
# ---------------------------
def tokenize_slides(text: str) -> List[str]:
    text = text.replace("\r", "")
    parts = re.split(r"(?mi)^\s*Slide\s*\d*\s*:\s*", text)
    if len(parts) > 1:
        blocks = []
        for part in parts[1:]:
            blocks.append(part.strip())
        return blocks
    return [blk.strip() for blk in re.split(r"\n\s*\n", text) if blk.strip()]

def parse_llm_response(response_text: str) -> List[Dict[str, Any]]:
    slides: List[Dict[str, Any]] = []
    blocks = tokenize_slides(response_text)

    for i, block in enumerate(blocks):
        lines = [l.rstrip() for l in block.split("\n") if l.strip()]
        if not lines:
            continue
        raw_title = lines[0]
        title = clean_inline(strip_slide_prefix(strip_leading_number(raw_title)))

        content_items: List[Dict[str, Any]] = []
        last_bullet_idx = None

        for line in lines[1:]:
            raw = line.strip()

            if raw.lower().startswith("subtitle:"):
                subtitle_text = clean_inline(raw.split(":", 1)[1])
                content_items.append({"type": "subtitle", "text": subtitle_text})
                continue

            if re.match(r"^[-*•]\s+", raw):
                body = strip_bullet_marker(raw)
                if ":" in body:
                    headline, detail = body.split(":", 1)
                    headline = clean_inline(strip_leading_number(headline))
                    detail = clean_inline(detail)
                    content_items.append({"type": "bullet", "title": headline, "detail": detail})
                else:
                    headline = clean_inline(strip_leading_number(body))
                    content_items.append({"type": "bullet", "title": headline, "detail": ""})
                last_bullet_idx = len(content_items) - 1
                continue

            if last_bullet_idx is not None and content_items[last_bullet_idx]["type"] == "bullet":
                extra = clean_inline(raw)
                if extra:
                    if content_items[last_bullet_idx]["detail"]:
                        content_items[last_bullet_idx]["detail"] += " " + extra
                    else:
                        content_items[last_bullet_idx]["detail"] = extra
                continue

            content_items.append({"type": "paragraph", "text": clean_inline(strip_leading_number(raw))})

        has_any = title or any(ci.get("text") or ci.get("title") for ci in content_items)
        if not has_any:
            continue

        slides.append({
            "title": title if title else f"Slide {i+1}",
            "content": content_items,
            "type": "title" if i == 0 else "content",
        })

    if not slides:
        slides = [{
            "title": "Presentation",
            "content": [{"type": "paragraph", "text": clean_inline(response_text)}],
            "type": "title",
        }]
    return slides

def api_key_required(text: str, guidance: str) -> bool:
    if guidance.strip():
        return True
    if len(text) > 3000:
        return True
    return False

# ---------------------------
# Slide creation
# ---------------------------
def _set_autofit(tf):
    # Shrink text to fit the textbox/placeholder
    try:
        tf.word_wrap = True
    except Exception:
        pass
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass

def create_slide(
    prs: Presentation,
    slide_data: Dict[str, Any],
    slide_index: int,
    *,
    force_layout0: bool
):
    """
    If force_layout0=True -> always use layout 0 (same as first slide) so color/background match.
    We then place our own text boxes for title/content to avoid missing placeholders.
    """
    layout = prs.slide_layouts[0] if (force_layout0 or slide_index == 0) else prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape
            slide.shapes._spTree.remove(sp._element)

    # ---- Title handling ----
    title_text = clean_inline(slide_data["title"])

    if slide.shapes.title:
        slide.shapes.title.text = title_text
        tf = slide.shapes.title.text_frame
        for p in tf.paragraphs:
            for run in p.runs:
                run.font.size = Pt(32 if slide_index == 0 else 28)
        if slide_index > 0:
            # Reposition title near top for consistent layout even if layout0
            slide.shapes.title.top = Inches(0.35)
            slide.shapes.title.left = Inches(0.6)
            slide.shapes.title.width = Inches(9.0)
            slide.shapes.title.height = Inches(0.9)
        _set_autofit(slide.shapes.title.text_frame)
    else:
        # If there is no title placeholder, create one
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(9.0), Inches(0.9))
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        _set_autofit(tf)

    # ---- Slide 1 subtitle handling ----
    if slide_index == 0:
        subtitle_items = [c for c in slide_data["content"] if c["type"] == "subtitle"]
        if subtitle_items:
            sub_text = subtitle_items[0]["text"]
            subtitle_shape = None
            try:
                if len(slide.placeholders) > 1 and slide.placeholders[1].has_text_frame:
                    subtitle_shape = slide.placeholders[1]
            except Exception:
                subtitle_shape = None

            if subtitle_shape is None:
                subtitle_shape = slide.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(8.0), Inches(1.2))

            tf = subtitle_shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = sub_text
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(20)
            _set_autofit(tf)
        return slide

    # ---- Slides 2+ content area ----
    # We will always create our own content textbox so every slide looks identical across templates.
    content_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(8.2), Inches(4.3))
    tf = content_box.text_frame
    tf.clear()
    _set_autofit(tf)

    first_para = True

    def add_para(text: str, level: int, *, bold=False, size=Pt(18), color: RGBColor | None = None):
        nonlocal first_para
        p = tf.paragraphs[0] if first_para else tf.add_paragraph()
        first_para = False
        p.text = text
        p.level = level
        p.font.bold = bold
        p.font.size = size
        if color:
            p.font.color.rgb = color
        return p

    for item in slide_data["content"]:
        if item["type"] == "bullet":
            head = item["title"]
            det  = item.get("detail", "")
            add_para(head, 0, bold=True, size=Pt(22))
            if det:
                add_para(det, 1, size=Pt(18), color=RGBColor(80, 80, 80))
        elif item["type"] == "paragraph":
            add_para(item["text"], 0, size=Pt(20))

    return slide

# ---------------------------
# Main pipeline
# ---------------------------
async def structured_markdown_to_slides(
    prs: Presentation,
    text_input: str,
    guidance: str,
    llm_provider: str,
    api_key: str,
    *,
    template_selected: bool
) -> Presentation:
    require_api = api_key_required(text_input, guidance)

    if require_api:
        if not api_key.strip():
            raise HTTPException(status_code=400, detail="API key required: text is too large or guidance was provided.")
        llm_response = await call_llm(llm_provider, api_key, text_input, guidance)
    else:
        llm_response = text_input

    slides_data = parse_llm_response(llm_response)

    # If template was chosen/uploaded, force every slide to use layout 0 so the background/design matches slide 1
    for i, slide_data in enumerate(slides_data):
        create_slide(prs, slide_data, i, force_layout0=template_selected)
    return prs

# ---------------------------
# API
# ---------------------------
@app.post("/api/convert")
async def convert(
    content: str = Form(None),
    markdown_file: UploadFile = File(None),
    template_file: UploadFile = File(None),
    template_id: str = Form(None),
    guidance: str = Form(""),
    llm_provider: str = Form("openai"),
    api_key: str = Form(""),
):
    # Handle content from either text box or markdown file
    if (not content or not content.strip()) and not markdown_file:
        raise HTTPException(status_code=400, detail="Please provide content in the text area or upload a Markdown file")

    if markdown_file:
        file_bytes = await markdown_file.read()
        content = file_bytes.decode("utf-8")

    prs: Presentation
    template_selected = False

    # Load template (pptx file or builtin template)
    if template_file:
        tmp_template = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        with open(tmp_template.name, "wb") as buffer:
            shutil.copyfileobj(template_file.file, buffer)
        prs = Presentation(tmp_template.name)
        template_selected = True
    elif template_id and template_id in BUILTIN_TEMPLATES and os.path.exists(BUILTIN_TEMPLATES[template_id]):
        prs = Presentation(BUILTIN_TEMPLATES[template_id])
        template_selected = True
    else:
        prs = Presentation()

    # Remove the initial default slide if present (many templates include one blank/title slide)
    try:
        if len(prs.slides) > 0:
            r_id = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(r_id)
            del prs.slides._sldIdLst[0]
    except Exception:
        pass

    # Convert content → slides (force layout0 if a template is selected so all slides match)
    prs = await structured_markdown_to_slides(
        prs, content, guidance, llm_provider, api_key, template_selected=template_selected
    )

    # Save and return file
    output_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(output_file.name)
    return FileResponse(
        output_file.name,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="slides.pptx",
    )

@app.get("/health")
def health():
    return {"ok": True}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("app:app", host="0.0.0.0", port=int(os.environ.get("PORT", "8001")), reload=False)
