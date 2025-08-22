import re, io
from dataclasses import dataclass, field
from typing import List, Optional

@dataclass
class ImgSpec:
    path: Optional[str] = None
    inline_bytes: Optional[bytes] = None

@dataclass
class SlideSpec:
    kind: str = "content"        # "title" or "content"
    title: str = ""
    text: str = ""
    bullets: List[str] = field(default_factory=list)
    levels: List[int] = field(default_factory=list)
    images: List[ImgSpec] = field(default_factory=list)
    table: List[List[str]] = field(default_factory=list)
    is_code: bool = False

MD_HR = re.compile(r"^\s*---\s*$")
MD_TITLE = re.compile(r"^\s*#\s+(.+)$")
MD_SUBTITLE = re.compile(r"^\s*##\s+(.+)$")
MD_BULLET = re.compile(r"^\s*[-*+]\s+(.+)$")
MD_NUM = re.compile(r"^\s*\d+\.\s+(.+)$")
MD_IMG = re.compile(r"!\[[^\]]*\]\(([^)]+)\)")
MD_TABLE_DIV = re.compile(r"^\s*\|")
MD_CODE_FENCE = re.compile(r"^\s*```")

from pptx.util import Pt, Inches

def add_slide(prs, title_text, body_lines, is_title_slide=False):
    if is_title_slide:
        slide_layout = prs.slide_layouts[0]
    else:
        slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(slide_layout)

    if slide.shapes.title:
        slide.shapes.title.text = title_text

    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = textbox.text_frame
    for line in body_lines:
        p = tf.add_paragraph()
        if line.startswith("CODE:"):
            p.text = line.replace("CODE:", "")
            p.font.name = "Courier New"
            p.font.size = Pt(16)
        elif line.startswith("• "):
            p.text = line[2:]
            p.level = 0
            p.font.size = Pt(18)
        else:
            p.text = line
            p.font.size = Pt(18)

def parse_markdown_to_slides(md: str, max_bullets_per_slide: int = 8) -> List[SlideSpec]:
    lines = md.splitlines()
    slides: List[SlideSpec] = []
    cur = SlideSpec(kind="content")
    in_code = False
    table_buf: List[str] = []

    def flush_table():
        nonlocal table_buf, cur
        if not table_buf:
            return
        rows = []
        for row in table_buf:
            parts = [c.strip() for c in row.strip().strip("|").split("|")]
            rows.append(parts)
        cur.table = rows
        table_buf = []

    def start_new_slide():
        nonlocal cur
        if cur.title or cur.text or cur.bullets or cur.images or cur.table:
            slides.append(cur)
        cur = SlideSpec(kind="content")

    for raw in lines:
        line = raw.rstrip("\n")

        # Code fences
        if MD_CODE_FENCE.match(line):
            in_code = not in_code
            if not in_code:
                cur.is_code = True
            continue
        if in_code:
            cur.text += (line + "\n")
            continue

        # Horizontal rule -> new slide
        if MD_HR.match(line):
            flush_table()
            start_new_slide()
            continue

        # Title
        m = MD_TITLE.match(line)
        if m:
            flush_table()
            start_new_slide()
            cur.kind = "title"
            cur.title = m.group(1).strip()
            continue

        # Subtitle as content slide title
        m = MD_SUBTITLE.match(line)
        if m:
            flush_table()
            start_new_slide()
            cur.kind = "content"
            cur.title = m.group(1).strip()
            continue

        # Bullets
        m = MD_BULLET.match(line) or MD_NUM.match(line)
        if m:
            flush_table()
            text = m.group(1).strip()
            # sub-bullet level via leading spaces
            indent = len(raw) - len(raw.lstrip(" "))
            level = min(indent // 2, 5)
            if len(cur.bullets) >= max_bullets_per_slide:
                # spill to next slide
                slides.append(cur)
                cur = SlideSpec(kind="content", title=cur.title)
            cur.bullets.append(text)
            cur.levels.append(level)
            continue

        # Images
        m = MD_IMG.search(line)
        if m:
            flush_table()
            path = m.group(1).strip()
            cur.images.append(ImgSpec(path=path))
            continue

        # Markdown table
        if MD_TABLE_DIV.match(line):
            table_buf.append(line)
            continue

        # Plain text → append to body
        if line.strip():
            cur.text += (line + "\n")

    flush_table()
    if cur.title or cur.text or cur.bullets or cur.images or cur.table:
        slides.append(cur)

    # Post: title-only slide gets kind=title
    for s in slides:
        if s.title and not (s.text or s.bullets or s.images or s.table):
            s.kind = "title"

    # Label continuation slides
    titles_seen = {}
    for s in slides:
        if s.title:
            titles_seen.setdefault(s.title, 0)
            if titles_seen[s.title] > 0:
                s.title = f"{s.title} (cont.)"
            titles_seen[s.title] += 1

    return slides